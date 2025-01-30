# requirements.txt
# pip install PyPDF2 python-docx python-pptx pytesseract transformers torch langchain nltk spacy textstat

####################################
### 1. Text Extraction Module ###
####################################

import re
from pathlib import Path
import PyPDF2
import docx
import pptx
import pytesseract
from PIL import Image
import io

class TextExtractor:
    def __init__(self):
        self.image_extensions = ['.png', '.jpg', '.jpeg']
    
    def extract(self, file_path):
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pdf':
            return self._extract_pdf(file_path)
        elif ext == '.docx':
            return self._extract_docx(file_path)
        elif ext == '.pptx':
            return self._extract_pptx(file_path)
        else:
            raise ValueError("Unsupported file type")

    def _extract_pdf(self, path):
        text = ""
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                if '/Image' in page['/Resources']:  # OCR for scanned PDFs
                    img = Image.open(io.BytesIO(page['/Resources']['/XObject']['/Im0'].get_data()))
                    text += pytesseract.image_to_string(img)
                else:
                    text += page.extract_text()
        return text
    
    def _extract_docx(self, path):
        doc = docx.Document(path)
        return '\n'.join([para.text for para in doc.paragraphs])
    
    def _extract_pptx(self, path):
        prs = pptx.Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

####################################
### 2. Text Processing & EDA ###
####################################

from langchain.text_splitter import RecursiveCharacterTextSplitter
import spacy
from nltk import sent_tokenize
import textstat

nlp = spacy.load("en_core_web_sm")

def process_text(text, chunk_size=512):
    # Clean text
    text = re.sub(r'\n{3,}', '\n\n', text)  # Remove excessive newlines
    
    # Split into context chunks
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=50
    )
    chunks = splitter.split_text(text)
    
    # Basic EDA
    readability = textstat.flesch_kincaid_grade(text)
    entities = [ent.text for ent in nlp(text).ents if ent.label_ in ['PERSON', 'ORG', 'LOC']]
    
    return {
        'chunks': chunks,
        'readability': readability,
        'key_entities': entities[:5]  # Top 5 entities
    }

####################################
### 3. Question Generation Model ###
####################################

from transformers import T5ForConditionedGeneration, T5Tokenizer
import torch

class QuizGenerator:
    def __init__(self, model_name="t5-base"):
        self.tokenizer = T5Tokenizer.from_pretrained(model_name)
        self.model = T5ForConditionedGeneration.from_pretrained(model_name)
        
        # Load fine-tuned weights (hypothetical path)
        # self.model.load_state_dict(torch.load("quiz_model.pth"))
    
    def generate(self, context, q_type="multiple_choice"):
        prompt = f"Generate {q_type} question: {context}"
        
        inputs = self.tokenizer(prompt, return_tensors="pt", max_length=512, truncation=True)
        outputs = self.model.generate(**inputs, max_length=200)
        
        return self.tokenizer.decode(outputs[0], skip_special_tokens=True)
    
    def batch_generate(self, chunks):
        questions = []
        for chunk in chunks:
            for q_type in ['multiple_choice', 'true_false', 'fill_blank']:
                question = self.generate(chunk, q_type)
                questions.append({
                    'question': question,
                    'context': chunk,
                    'type': q_type
                })
        return questions

####################################
### 4. Validation & Post-Processing ###
####################################

from transformers import pipeline

class QuizValidator:
    def __init__(self):
        self.qa_pipeline = pipeline("question-answering", model="deepset/roberta-base-squad2")
    
    def validate_answer(self, question, context, answer):
        result = self.qa_pipeline(question=question, context=context)
        return result['answer'] == answer
    
    def check_readability(self, text, max_grade=3):
        score = textstat.flesch_kincaid_grade(text)
        return score <= max_grade

####################################
### 5. End-to-End Pipeline ###
####################################

class QuizPipeline:
    def __init__(self):
        self.extractor = TextExtractor()
        self.generator = QuizGenerator()
        self.validator = QuizValidator()
    
    def run(self, file_path):
        # Step 1: Extract text
        raw_text = self.extractor.extract(file_path)
        
        # Step 2: Process text
        processed = process_text(raw_text)
        
        # Step 3: Generate questions
        questions = self.generator.batch_generate(processed['chunks'])
        
        # Step 4: Validate
        valid_questions = []
        for q in questions:
            if self.validator.check_readability(q['question']):
                valid_questions.append(q)
        
        return {
            'metadata': {
                'readability_score': processed['readability'],
                'key_entities': processed['key_entities']
            },
            'questions': valid_questions[:10]  # Return top 10
        }

# Example Usage
if __name__ == "__main__":
    pipeline = QuizPipeline()
    result = pipeline.run("sample_lesson.pdf")
    print(f"Generated {len(result['questions'])} questions:")
    for q in result['questions']:
        print(f"[{q['type'].upper()}] {q['question']}")
